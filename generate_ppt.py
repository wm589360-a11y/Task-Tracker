import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation()

# Define Slide Content
slides_data = [
    {
        "title": "TaskTracker",
        "subtitle": "Streamlining Project Management\nEmpowering teams to collaborate, track progress, and deliver results.",
        "type": "title"
    },
    {
        "title": "The Challenges We Face",
        "points": [
            "Scattered communication across emails and multiple chat apps.",
            "Unclear responsibilities leading to missed deadlines.",
            "Lack of real-time visibility into overall project progress."
        ]
    },
    {
        "title": "Proposed Solution: Overview",
        "points": [
            "A centralized web-based platform for all task management.",
            "Unified communication, file sharing, and status tracking in one place.",
            "Intuitive dashboards to visualize workload and progress."
        ]
    },
    {
        "title": "Project Goals & Objectives",
        "points": [
            "Efficiency: Improve team productivity by reducing time spent searching for task details.",
            "Accountability: Clearly define roles and assign specific tasks to individuals.",
            "Transparency: Provide real-time progress tracking for managers and clients.",
            "Scalability: Build a system that grows alongside the organization."
        ]
    },
    {
        "title": "Server Requirements",
        "points": [
            "Hosting: Cloud-based VPS (e.g., AWS, DigitalOcean) or standard web hosting.",
            "Web Server: Apache or Nginx.",
            "Hardware (Minimum): 2GB RAM, 2 CPU Cores, 20GB SSD Storage.",
            "OS: Linux (Ubuntu/CentOS)."
        ]
    },
    {
        "title": "Technology Stack",
        "points": [
            "Frontend: HTML5, CSS3, JavaScript (React / Vanilla JS).",
            "Backend: PHP / Node.js.",
            "Database: MySQL / PostgreSQL.",
            "Version Control: Git & GitHub."
        ]
    },
    {
        "title": "End User Interface",
        "points": [
            "Personal Dashboard: Summary of assigned tasks and upcoming deadlines.",
            "Kanban Board / List View: Easy drag-and-drop task management.",
            "Task Details: Spaces for comments, attachments, and status updates.",
            "Notifications: Real-time alerts for updates and mentions."
        ]
    },
    {
        "title": "Admin Dashboard",
        "points": [
            "User Management: Invite, suspend, or remove team members.",
            "Role-Based Access: Assign Admin, Manager, or Employee permissions.",
            "Global Overview: Monitor all active projects across the company.",
            "Analytics: Generate reports on team performance and bottlenecks."
        ]
    },
    {
        "title": "Database Design",
        "points": [
            "Relational Structure: Ensuring data integrity and reducing redundancy.",
            "Core Tables: Users, Projects, Tasks, Comments, Attachments.",
            "Security: Hashed passwords and secure session management."
        ]
    },
    {
        "title": "Entity-Relationship (ER) Diagram",
        "points": [
            "Key Relationships:",
            "  • 1 User -> Many Tasks (One-to-Many)",
            "  • 1 Project -> Many Tasks (One-to-Many)",
            "  • 1 Task -> Many Comments (One-to-Many)",
            "\n[Please insert your ER Diagram Image Here]"
        ]
    },
    {
        "title": "Planning & Data Flow",
        "points": [
            "DFD (Data Flow Diagram): Shows how user inputs travel to the database.",
            "Gantt Chart: Illustrates our development timeline and milestones.",
            "PERT Chart: Highlights task dependencies during development.",
            "\n[Please insert visual diagrams Here]"
        ]
    },
    {
        "title": "Modules & Features",
        "points": [
            "Authentication Module: Secure Login/Signup & Password Reset.",
            "Task Management Module: Create, Read, Update, Delete (CRUD) tasks.",
            "Collaboration Module: File sharing and real-time commenting.",
            "Search & Filter Module: Quickly find tasks by name, date, or assignee."
        ]
    },
    {
        "title": "Challenges & Learning",
        "points": [
            "Challenges: Handling database relationships, creating a responsive UI, and managing state/data flow.",
            "Learnings: Deepened understanding of relational databases, improved UI/UX design skills, and the importance of testing."
        ]
    },
    {
        "title": "References & Resources",
        "points": [
            "Official Documentation (React, Node.js, PHP, MySQL, etc.)",
            "Stack Overflow / Developer Forums",
            "UI/UX Inspiration: Dribbble / Behance",
            "Open Source Libraries and Frameworks used."
        ]
    }
]

# Layouts
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]

for slide_data in slides_data:
    if slide_data.get("type") == "title":
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = slide_data["title"]
        subtitle.text = slide_data["subtitle"]
    else:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = slide_data["title"]
        
        tf = body_shape.text_frame
        for i, point in enumerate(slide_data["points"]):
            if i == 0:
                tf.text = point
            else:
                p = tf.add_paragraph()
                p.text = point
                
                if point.startswith("  "):
                    p.level = 1
                else:
                    p.level = 0

prs.save("TaskTracker_Presentation.pptx")
print("Presentation generated successfully as TaskTracker_Presentation.pptx")
